from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "interactive-slides.html"
OUTPUT = ROOT / "政好_可信政府AI_Agent_Demo.pptx"
SLIDE_COUNT = 14
WIDTH = 1920
HEIGHT = 1080


def find_chrome() -> Path:
    candidates = (
        Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
        Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
        Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    )
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("Chrome or Edge was not found.")


def capture_slides(chrome: Path, output_dir: Path) -> list[Path]:
    screenshots = []
    profile_dir = output_dir / "browser-profile"
    for slide_number in range(1, SLIDE_COUNT + 1):
        screenshot = output_dir / f"slide-{slide_number:02d}.png"
        url = f"{SOURCE.as_uri()}?export=ppt#{slide_number}"
        subprocess.run(
            [
                str(chrome),
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                "--allow-file-access-from-files",
                f"--user-data-dir={profile_dir}",
                f"--window-size={WIDTH},{HEIGHT}",
                "--force-device-scale-factor=1",
                "--run-all-compositor-stages-before-draw",
                "--virtual-time-budget=2500",
                f"--screenshot={screenshot}",
                url,
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        with Image.open(screenshot) as image:
            if image.size != (WIDTH, HEIGHT):
                raise RuntimeError(
                    f"Slide {slide_number} is {image.size}, expected {(WIDTH, HEIGHT)}."
                )
        screenshots.append(screenshot)
    return screenshots


def create_presentation(screenshots: list[Path]) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for screenshot in screenshots:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(screenshot),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    presentation.core_properties.title = "政好：可信政府 AI Agent"
    presentation.core_properties.subject = "可信 AI 黑客松 Demo 簡報"
    presentation.core_properties.author = "Ali、慶霖、Steve"
    presentation.save(OUTPUT)


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    chrome = find_chrome()
    temp_dir = Path(tempfile.mkdtemp(prefix="trusted-agent-ppt-"))
    try:
        screenshots = capture_slides(chrome, temp_dir)
        create_presentation(screenshots)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    print(OUTPUT)


if __name__ == "__main__":
    main()
